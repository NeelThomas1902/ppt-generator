"""Vision service – analyse existing PowerPoint files to extract structure."""

from __future__ import annotations

import json
from typing import Any, Dict, List

from pptx import Presentation
from pptx.util import Pt


class VisionService:
    """Extracts layout and content information from an existing .pptx file."""

    def analyse(self, file_path: str) -> Dict[str, Any]:
        """Return a structured representation of *file_path*."""
        prs = Presentation(file_path)

        slides: List[Dict[str, Any]] = []
        for slide in prs.slides:
            slide_info: Dict[str, Any] = {
                "layout": slide.slide_layout.name,
                "shapes": [],
            }
            for shape in slide.shapes:
                shape_info: Dict[str, Any] = {
                    "name": shape.name,
                    "shape_type": str(shape.shape_type),
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                }
                if shape.has_text_frame:
                    shape_info["text"] = shape.text_frame.text
                    shape_info["paragraphs"] = [
                        {
                            "text": para.text,
                            "level": para.level,
                            "font_size": (
                                para.runs[0].font.size.pt
                                if para.runs and para.runs[0].font.size
                                else None
                            ),
                        }
                        for para in shape.text_frame.paragraphs
                    ]
                slide_info["shapes"].append(shape_info)
            slides.append(slide_info)

        return {
            "slide_count": len(slides),
            "slide_width": prs.slide_width,
            "slide_height": prs.slide_height,
            "slides": slides,
        }

    def extract_theme_colors(self, file_path: str) -> List[str]:
        """Return a list of hex color strings used in the presentation."""
        prs = Presentation(file_path)
        colors: List[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.color and run.font.color.rgb:
                                hex_color = str(run.font.color.rgb)
                                if hex_color not in colors:
                                    colors.append(hex_color)
        return colors
