"""PowerPoint generation service."""

from __future__ import annotations

import os
import uuid
from datetime import datetime, timezone
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.api.schemas import PresentationResponse
from app.config import settings
from app.services.llm_types import LLMClient


_LAYOUT_TITLE = 0
_LAYOUT_TITLE_CONTENT = 1
_LAYOUT_BLANK = 6


class PPTGenerator:
    """Creates .pptx files from LLM-generated slide content."""

    def __init__(self, llm_service: LLMClient) -> None:
        self._llm = llm_service

    async def generate(
        self,
        prompt: str,
        template_id: Optional[str] = None,
        slide_count: int = 8,
        theme: Optional[str] = None,
        template_definition: Optional[Dict[str, Any]] = None,
    ) -> PresentationResponse:
        """Generate a presentation and return metadata."""
        content = await self._llm.generate_presentation_content(
            prompt=prompt,
            slide_count=slide_count,
            theme=theme,
        )

        prs = None
        if template_definition:
            prs = self._build_from_template(template_definition, content)

        if prs is None:
            prs = Presentation()
            self._apply_theme(prs, theme)

            # Title slide
            title_layout = prs.slide_layouts[_LAYOUT_TITLE]
            title_slide = prs.slides.add_slide(title_layout)
            title_slide.shapes.title.text = content.get("title", "Presentation")

            # Content slides
            for slide_data in content.get("slides", []):
                layout = prs.slide_layouts[_LAYOUT_TITLE_CONTENT]
                slide = prs.slides.add_slide(layout)

                if slide.shapes.title:
                    slide.shapes.title.text = slide_data.get("title", "")

                body = slide.placeholders[1]
                tf = body.text_frame
                tf.clear()

                for i, bullet in enumerate(slide_data.get("content", [])):
                    if i == 0:
                        tf.paragraphs[0].text = bullet
                    else:
                        p = tf.add_paragraph()
                        p.text = bullet
                        p.level = 0

                if slide_data.get("notes"):
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = slide_data["notes"]

        # Persist
        os.makedirs(settings.generated_dir, exist_ok=True)
        filename = f"{uuid.uuid4()}.pptx"
        output_path = os.path.join(settings.generated_dir, filename)
        prs.save(output_path)

        return PresentationResponse(
            id=str(uuid.uuid4()),
            filename=filename,
            download_url=f"/generated/{filename}",
            slide_count=len(prs.slides),
            created_at=datetime.now(timezone.utc).isoformat(),
            prompt=prompt,
        )

    # ── Internal helpers ───────────────────────────────────────────────────────

    def _apply_theme(self, prs: Presentation, theme: Optional[str]) -> None:
        """Apply very basic theme colours to the slide master."""
        if theme == "corporate":
            # Dark blue background for title layouts (cosmetic only)
            pass
        # Default: no changes – use python-pptx built-in defaults

    def _build_from_template(
        self,
        template_definition: Dict[str, Any],
        content: Dict[str, Any],
    ) -> Presentation:
        template_file = template_definition.get("template_file")
        if template_file:
            return self._build_from_potx(template_file, content, template_definition)

        prs = Presentation()
        self._apply_slide_size(prs, template_definition)

        slides_def = template_definition.get("slides", [])
        title_template = self._pick_template_slide(slides_def, need_title=True, need_body=False)
        body_template = self._pick_template_slide(slides_def, need_title=True, need_body=True)
        if not body_template:
            body_template = self._pick_template_slide(slides_def, need_title=False, need_body=True)

        title_slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_BLANK])
        if title_template:
            self._render_template_slide(
                title_slide,
                title_template,
                {"title": content.get("title", "Presentation"), "content": []},
            )
        else:
            title_slide.shapes.title.text = content.get("title", "Presentation")

        for slide_data in content.get("slides", []):
            slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_BLANK])
            if body_template:
                self._render_template_slide(slide, body_template, slide_data)
            else:
                self._render_simple_slide(slide, slide_data)

        return prs

    def _build_from_potx(
        self,
        template_file: str,
        content: Dict[str, Any],
        template_definition: Dict[str, Any],
    ) -> Presentation:
        template_path = self._resolve_template_path(template_file)
        prs = Presentation(template_path)
        self._clear_slides(prs)

        title_layout = self._find_layout_with_placeholders(
            prs,
            {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE},
        )
        body_layout = self._find_layout_with_placeholders(
            prs,
            {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT},
        )
        if body_layout is None:
            body_layout = title_layout

        slides_def = template_definition.get("slides", [])
        title_template = self._pick_template_slide(slides_def, need_title=True, need_body=False)
        body_template = self._pick_template_slide(slides_def, need_title=True, need_body=True)
        if not body_template:
            body_template = self._pick_template_slide(slides_def, need_title=False, need_body=True)

        title_idx = self._first_placeholder_index(title_template, "title")
        body_indices = self._placeholder_indices(body_template, "body")

        if title_layout:
            title_slide = prs.slides.add_slide(title_layout)
            title_placeholder = None
            if title_idx is not None:
                title_placeholder = self._find_placeholder_by_idx(title_slide, title_idx)
            if not title_placeholder:
                title_placeholder = self._find_placeholder(
                    title_slide,
                    {
                        PP_PLACEHOLDER.TITLE,
                        PP_PLACEHOLDER.CENTER_TITLE,
                    },
                )
            if not title_placeholder:
                title_placeholder = self._find_text_shape(
                    title_slide,
                    exclude_types={PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT},
                )
            if title_placeholder and title_placeholder.has_text_frame:
                title_placeholder.text_frame.clear()
                title_placeholder.text_frame.text = content.get("title", "Presentation")
            else:
                self._add_fallback_textbox(
                    title_slide,
                    content.get("title", "Presentation"),
                    top_ratio=0.18,
                    height_ratio=0.2,
                    font_size=Pt(36),
                )

        for slide_data in content.get("slides", []):
            layout = body_layout or prs.slide_layouts[_LAYOUT_TITLE_CONTENT]
            slide = prs.slides.add_slide(layout)
            title_placeholder = None
            if title_idx is not None:
                title_placeholder = self._find_placeholder_by_idx(slide, title_idx)
            if not title_placeholder:
                title_placeholder = self._find_placeholder(
                    slide,
                    {
                        PP_PLACEHOLDER.TITLE,
                        PP_PLACEHOLDER.CENTER_TITLE,
                    },
                )
            if title_placeholder and title_placeholder.has_text_frame:
                title_placeholder.text_frame.clear()
                title_placeholder.text_frame.text = slide_data.get("title", "")

            body = None
            for idx in body_indices:
                body = self._find_placeholder_by_idx(slide, idx)
                if body is not None:
                    break
            if not body:
                body = self._find_placeholder(
                    slide,
                    {
                        PP_PLACEHOLDER.BODY,
                        PP_PLACEHOLDER.OBJECT,
                    },
                )
            if not body:
                body = self._find_text_shape(
                    slide,
                    exclude_types={
                        PP_PLACEHOLDER.TITLE,
                        PP_PLACEHOLDER.CENTER_TITLE,
                    },
                )
            if body and body.has_text_frame:
                tf = body.text_frame
                tf.clear()
                for i, bullet in enumerate(slide_data.get("content", [])):
                    if i == 0:
                        tf.paragraphs[0].text = bullet
                    else:
                        p = tf.add_paragraph()
                        p.text = bullet
                        p.level = 0
            elif slide_data.get("content"):
                self._add_fallback_textbox(
                    slide,
                    "\n".join(slide_data.get("content", [])),
                    top_ratio=0.28,
                    height_ratio=0.55,
                    font_size=Pt(18),
                )

        return prs

    def _resolve_template_path(self, template_file: str) -> str:
        if os.path.isabs(template_file):
            return template_file
        base_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
        return os.path.join(base_dir, template_file)

    def _clear_slides(self, prs: Presentation) -> None:
        slide_ids = list(prs.slides._sldIdLst)
        for slide_id in slide_ids:
            r_id = slide_id.rId
            prs.slides._sldIdLst.remove(slide_id)
            prs.part.drop_rel(r_id)

    def _find_placeholder(self, slide, types: set[PP_PLACEHOLDER]):
        for shape in slide.placeholders:
            if not shape.is_placeholder:
                continue
            if shape.placeholder_format.type in types:
                return shape
        return None

    def _find_text_shape(
        self,
        slide,
        exclude_types: set[PP_PLACEHOLDER] | None = None,
    ):
        exclude_types = exclude_types or set()
        for shape in slide.placeholders:
            if not shape.is_placeholder or not shape.has_text_frame:
                continue
            if shape.placeholder_format.type in exclude_types:
                continue
            return shape
        for shape in slide.shapes:
            if shape.has_text_frame:
                return shape
        return None

    def _find_placeholder_by_idx(self, slide, idx: int):
        for shape in slide.placeholders:
            if not shape.is_placeholder:
                continue
            if shape.placeholder_format.idx == idx:
                return shape
        return None

    def _find_layout_with_placeholders(
        self,
        prs: Presentation,
        types: set[PP_PLACEHOLDER],
    ):
        for layout in prs.slide_layouts:
            for shape in layout.placeholders:
                if shape.placeholder_format.type in types:
                    return layout
        return None

    def _add_fallback_textbox(
        self,
        slide,
        text: str,
        top_ratio: float,
        height_ratio: float,
        font_size: Pt,
    ) -> None:
        prs_width = slide.part.slide_layout.part.presentation.slide_width
        prs_height = slide.part.slide_layout.part.presentation.slide_height
        left = int(prs_width * 0.08)
        top = int(prs_height * top_ratio)
        box_width = int(prs_width * 0.84)
        box_height = int(prs_height * height_ratio)
        shape = slide.shapes.add_textbox(left, top, box_width, box_height)
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = font_size

    def _placeholder_indices(
        self,
        template_slide: Optional[Dict[str, Any]],
        placeholder_type: str,
    ) -> List[int]:
        indices: List[int] = []
        if not template_slide:
            return indices
        target = placeholder_type.lower()
        for shape in template_slide.get("shapes", []):
            if shape.get("placeholderType", "").lower() != target:
                continue
            idx = shape.get("placeholderIndex")
            if isinstance(idx, int):
                indices.append(idx)
        return indices

    def _first_placeholder_index(
        self,
        template_slide: Optional[Dict[str, Any]],
        placeholder_type: str,
    ) -> Optional[int]:
        indices = self._placeholder_indices(template_slide, placeholder_type)
        return indices[0] if indices else None

    def _apply_slide_size(self, prs: Presentation, template_definition: Dict[str, Any]) -> None:
        slide_size = template_definition.get("slideSize") or {}
        width = slide_size.get("width")
        height = slide_size.get("height")
        if width and height:
            prs.slide_width = Inches(width)
            prs.slide_height = Inches(height)

    def _pick_template_slide(
        self,
        slides_def: List[Dict[str, Any]],
        need_title: bool,
        need_body: bool,
    ) -> Optional[Dict[str, Any]]:
        for slide in slides_def:
            shapes = slide.get("shapes", [])
            has_title = any(s.get("placeholderType") == "title" for s in shapes)
            has_body = any(s.get("placeholderType") == "body" for s in shapes)
            if (not need_title or has_title) and (not need_body or has_body):
                return slide
        return slides_def[0] if slides_def else None

    def _render_simple_slide(self, slide, slide_data: Dict[str, Any]) -> None:
        title = slide.shapes.title
        if title:
            title.text = slide_data.get("title", "")
        body = slide.placeholders[1] if slide.placeholders else None
        if not body:
            return
        tf = body.text_frame
        tf.clear()
        for i, bullet in enumerate(slide_data.get("content", [])):
            if i == 0:
                tf.paragraphs[0].text = bullet
            else:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

    def _render_template_slide(
        self,
        slide,
        template_slide: Dict[str, Any],
        slide_data: Dict[str, Any],
    ) -> None:
        for shape_def in template_slide.get("shapes", []):
            box = self._shape_box(shape_def)
            if not box:
                continue

            if shape_def.get("type") == "picture":
                self._add_picture_placeholder(slide, shape_def, box)
                continue

            placeholder_type = shape_def.get("placeholderType")
            if placeholder_type == "title" and slide_data.get("title"):
                text_items = [{"text": slide_data.get("title", "")}]
            elif placeholder_type == "body" and slide_data.get("content"):
                text_items = [{"text": item, "level": 0} for item in slide_data.get("content", [])]
            else:
                text_items = shape_def.get("textContent")

            if text_items:
                self._add_text_shape(slide, shape_def, box, text_items)
            else:
                self._add_shape(slide, shape_def, box)

    def _add_text_shape(
        self,
        slide,
        shape_def: Dict[str, Any],
        box: Dict[str, float],
        text_items: List[Dict[str, Any]],
    ) -> None:
        shape = slide.shapes.add_textbox(
            Inches(box["x"]),
            Inches(box["y"]),
            Inches(box["width"]),
            Inches(box["height"]),
        )
        self._apply_fill(shape, shape_def)
        tf = shape.text_frame
        tf.word_wrap = True
        self._apply_text_items(tf, text_items)

    def _add_shape(self, slide, shape_def: Dict[str, Any], box: Dict[str, float]) -> None:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(box["x"]),
            Inches(box["y"]),
            Inches(box["width"]),
            Inches(box["height"]),
        )
        self._apply_fill(shape, shape_def)
        shape.line.fill.background()

    def _add_picture_placeholder(self, slide, shape_def: Dict[str, Any], box: Dict[str, float]) -> None:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(box["x"]),
            Inches(box["y"]),
            Inches(box["width"]),
            Inches(box["height"]),
        )
        shape.fill.background()
        shape.line.fill.background()

    def _shape_box(self, shape_def: Dict[str, Any]) -> Optional[Dict[str, float]]:
        pos = shape_def.get("position") or {}
        size = shape_def.get("size") or {}
        if "width" not in size or "height" not in size:
            return None
        return {
            "x": float(pos.get("x", 0.0)),
            "y": float(pos.get("y", 0.0)),
            "width": float(size.get("width", 1.0)),
            "height": float(size.get("height", 1.0)),
        }

    def _apply_text_items(self, text_frame, text_items: List[Dict[str, Any]]) -> None:
        text_frame.clear()
        for i, item in enumerate(text_items):
            paragraph = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            level = item.get("level")
            if level is not None:
                paragraph.level = int(level)
            alignment = item.get("alignment")
            if alignment == "ctr":
                paragraph.alignment = PP_ALIGN.CENTER
            elif alignment == "r":
                paragraph.alignment = PP_ALIGN.RIGHT
            elif alignment == "l":
                paragraph.alignment = PP_ALIGN.LEFT

            runs = item.get("runs")
            if runs:
                for run_def in runs:
                    run = paragraph.add_run()
                    run.text = run_def.get("text", "")
                    if "bold" in run_def:
                        run.font.bold = bool(run_def["bold"])
                    if "fontSize" in run_def:
                        run.font.size = Pt(float(run_def["fontSize"]))
                    color = self._hex_to_rgb(run_def.get("color"))
                    if color:
                        run.font.color.rgb = color
            else:
                paragraph.text = item.get("text", "")

    def _apply_fill(self, shape, shape_def: Dict[str, Any]) -> None:
        fill_color = shape_def.get("fillColor")
        rgb = self._hex_to_rgb(fill_color)
        if rgb:
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb

    def _hex_to_rgb(self, value: Optional[str]) -> Optional[RGBColor]:
        if not value:
            return None
        hex_value = value.lstrip("#")
        if len(hex_value) != 6:
            return None
        try:
            r = int(hex_value[0:2], 16)
            g = int(hex_value[2:4], 16)
            b = int(hex_value[4:6], 16)
        except ValueError:
            return None
        return RGBColor(r, g, b)
